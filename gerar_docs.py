import os
import json
import pandas as pd
from pptx import Presentation
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

# 1. Garantir que a pasta de documentos existe
os.makedirs('docs_auramotors', exist_ok=True)

# 2. PDF: Manual do Robô KR-500 (usando ReportLab)
pdf_path = 'docs_auramotors/manual_robo_kr500.pdf'
c = canvas.Canvas(pdf_path, pagesize=letter)
c.setFont("Helvetica-Bold", 16)
c.drawString(50, 750, "AuraMotors - Manual Técnico Robô de Solda KR-500")

c.setFont("Helvetica", 10)
c.drawString(50, 710, "Código de Erro E-102: Superaquecimento da junta articulada 3.")
c.drawString(50, 695, "Procedimento: Desligar o braço por 15 minutos e checar nível do fluido refrigerante.")
c.drawString(50, 665, "Manutenção Preventiva: Troca do óleo lubrificante a cada 500 horas de operação.")
c.save()

# 3. CSV: Tabela de Estoque de Insumos
df_estoque = pd.DataFrame({
    'SKU': ['SKU-8821', 'SKU-4412', 'SKU-9901'],
    'Item': ['Parafuso M8 Aço Inox', 'Sensor Indutivo 24V', 'Rolamento Axial 35mm'],
    'Quantidade': [1500, 12, 45],
    'Estoque_Minimo': [500, 20, 50],
    'Status': ['OK', 'REPOSIÇÃO CRÍTICA', 'REPOSIÇÃO NECESSÁRIA']
})
df_estoque.to_csv('docs_auramotors/estoque_insumos.csv', index=False)

# 4. Excel (.xlsx): Relatório Financeiro de Peças
with pd.ExcelWriter('docs_auramotors/relatorio_compras.xlsx', engine='openpyxl') as writer:
    df_estoque.to_excel(writer, sheet_name='Inventario', index=False)

# 5. JSON: SLA e Cadastro de Fornecedores
fornecedores = [
    {"id": "FORN-01", "nome": "Metalúrgica Precision", "categoria": "Fixadores", "sla_entrega_dias": 3},
    {"id": "FORN-02", "nome": "AutoSensors Brasil", "categoria": "Sensores", "sla_entrega_dias": 7}
]
with open('docs_auramotors/fornecedores.json', 'w', encoding='utf-8') as f:
    json.dump(fornecedores, f, ensure_ascii=False, indent=4)

# 6. Word (.docx): Procedimento Operacional Padrão (POP)
doc = Document()
doc.add_heading('POP-MN-042: Lubrificação de Prensas Hidráulicas', 0)
doc.add_paragraph('Objetivo: Padronizar o processo de manutenção preventiva nas prensas da Linha B.')
doc.add_paragraph('Frequência: Semanal. Graxa recomendada: ISO VG 220.')
doc.save('docs_auramotors/pop_lubrificacao.docx')

# 7. Markdown (.md): Guia Rápido de Diagnósticos
md_text = """# Guia Rápido de Diagnóstico - Linha de Montagem AuraMotors

## Ruídos e Vibrações na Esteira 02
1. Verificar se há obstrução no trilho guia.
2. Checar alinhamento do motor trifásico.
3. Se a vibração exceder 4.5 mm/s, acionar a equipe de manutenção preditiva imediatamente.
"""
with open('docs_auramotors/diagnostico_esteira.md', 'w', encoding='utf-8') as f:
    f.write(md_text)

# 8. PowerPoint (.pptx): Metas e OKRs
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AuraMotors - Metas de Manutenção Q3/Q4"
subtitle.text = "Meta 1: Reduzir tempo de paradas não planejadas (Downtime) em 15%.\nMeta 2: Manter 98% de precisão no estoque de SKUs críticos."
prs.save('docs_auramotors/metas_manutencao.pptx')

# 9. HTML: Comunicado Interno
html_pagina = """
<!DOCTYPE html>
<html>
<head><title>Política de Solicitação de Peças</title></head>
<body>
    <h1>Política Interna: Requisição de Peças no Chão de Fábrica</h1>
    <p>Todas as solicitações de peças com valor acima de R$ 2.000,00 exigem aprovação prévia do Supervisor de Turno via sistema ERP.</p>
</body>
</html>
"""
with open('docs_auramotors/politica_requisicao.html', 'w', encoding='utf-8') as f:
    f.write(html_pagina)

print("Base de conhecimento sintética da AuraMotors gerada com sucesso!")