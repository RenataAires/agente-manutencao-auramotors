import os
import json
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from dotenv import load_dotenv
from google import genai

# 1. Carrega a chave de forma segura
load_dotenv()

client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))

DOCS_DIR = "docs_auramotors"

def carregar_conteudo_documentos():
    """Percorre a pasta de documentos e extrai o texto de cada formato."""
    contexto_acumulado = []

    for root, _, files in os.walk(DOCS_DIR):
        for file in files:
            filepath = os.path.join(root, file)
            ext = file.split('.')[-1].lower()
            
            try:
                conteudo = ""
                # Lógica simples de extração baseada na extensão do arquivo
                if ext == "pdf":
                    reader = PdfReader(filepath)
                    for page in reader.pages:
                        conteudo += page.extract_text() or ""
                elif ext in ["csv", "xlsx"]:
                    df = pd.read_csv(filepath) if ext == "csv" else pd.read_excel(filepath)
                    conteudo = df.to_string()
                elif ext == "json":
                    with open(filepath, 'r', encoding='utf-8') as f:
                        conteudo = json.dumps(json.load(f), ensure_ascii=False)
                elif ext == "docx":
                    doc = DocxDocument(filepath)
                    conteudo = "\n".join([p.text for p in doc.paragraphs])
                elif ext == "pptx":
                    prs = Presentation(filepath)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                conteudo += shape.text + "\n"
                elif ext in ["md", "html"]:
                    with open(filepath, 'r', encoding='utf-8') as f:
                        conteudo = f.read()

                # Se conseguiu ler algo, adiciona ao "pacote" de contexto
                if conteudo.strip():
                    contexto_acumulado.append(f"--- INÍCIO: {file} ---\n{conteudo}\n--- FIM: {file} ---\n")
            
            except Exception as e:
                print(f"Erro ao ler {file}: {e}")

    return "\n\n".join(contexto_acumulado)

def perguntar_ao_agente(pergunta: str) -> str:
    """Cria a sessão de chat estruturada e envia a pergunta para a IA."""
    contexto = carregar_conteudo_documentos()
    
    system_instruction = f"""
    Você é o Agente de IA Corporativo da AuraMotors Componentes Automotivos S.A.
    Sua missão é responder com precisão, clareza e tom profissional às dúvidas dos colaboradores.
    Utilize EXCLUSIVAMENTE a base de conhecimento abaixo para responder. 
    
    BASE DE CONHECIMENTO:
    {contexto}
    """
    
    chat = client.chats.create(
        model="gemini-3.5-flash",
        config={"system_instruction": system_instruction}
    )
    
    response = chat.send_message(pergunta)
    return response.text

if __name__ == "__main__":
    # Teste de execução no terminal
    pergunta_teste = "Qual é o procedimento para o erro E-102 no robô KR-500?"
    print(f"Pergunta: {pergunta_teste}\n")
    print("Consultando a base de dados...\n")
    
    resposta = perguntar_ao_agente(pergunta_teste)
    print(f"Resposta do Agente:\n{resposta}")